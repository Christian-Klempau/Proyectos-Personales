from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
import os

path = os.path.join('ppt_test','new_ppt.pptx')
prs = Presentation(path)

# COLUMN_CLUSTERED (51)
# COLUMN_CLUSTERED (51)
# DOUGHNUT (-4120)
# DOUGHNUT (-4120)
# COLUMN_CLUSTERED (51)
datos = open("datos.csv", "r", encoding="UTF-8")
lista = []
for linea in datos:
    linea = linea.strip()
    linea = linea.split(",")
    for i in range(1, len(linea)):
        linea[i] = linea[i].split(";")
        lis = []
        for elem in linea[i]:
            lis.append(float(elem))
        linea[i] = lis
    lista.append(linea)
print(lista)
cnt = 0
for slide in prs.slides:
    idx = 1
    for shape in slide.shapes:
        if shape.has_text_frame:
            # ! N = x
            if shape.text_frame.text.lower().find("cultura, relaciones corporativas,") != -1:
                shape.text_frame.text = lista[cnt][0]
            pass
        if shape.has_chart:
            chart = shape.chart
            if str(chart.chart_type).find("COLUMN_CLUSTERED") != -1 and idx < 5:
                chart_data = CategoryChartData()
                chart_data.categories = ["Chart"]
                if idx == 2:
                    chart_data.add_series('Series 1', (lista[cnt][1][0], None))
                if idx == 1:
                    chart_data.add_series('Series 1', (lista[cnt][3][0], None))
                chart.replace_data(chart_data)
                #print(chart.chart_type, idx)
                idx += 1
            elif str(chart.chart_type).find("DOUGHNUT") != -1 and idx < 5:
                chart_data = ChartData()
                chart_data.categories = ["Chart"]
                if idx == 3:
                    chart_data.add_series(None, (lista[cnt][2][0]/100, lista[cnt][2][1]/100))
                if idx == 4:
                    chart_data.add_series(None, (lista[cnt][4][0]/100, lista[cnt][4][1]/100, lista[cnt][4][2]/100))
                chart.replace_data(chart_data)
                #print(chart.chart_type, idx)


                #print(chart.chart_type)
                idx += 1
            #print(idx)
    print(cnt)
    cnt += 1
    if cnt >= len(lista) - 1:
        cnt = len(lista) - 1
    
    #print("-------------")
prs.save(path)
print("DONE")